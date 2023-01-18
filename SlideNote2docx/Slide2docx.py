# !pip install python-pptx python-docx

from pptx import Presentation
from docx import Document
from glob import glob
from docx.shared import Mm

_PPTX = 'SlideNote2docx\\ex.pptx'
_DOCX = 'SlideNote2docx\\ex.docx'
_IMGS = 'SlideNote2docx\\ex\\*'

size_settings = {
    "page_height": Mm(297),
    "page_width": Mm(210),
    "left_margin": Mm(12.5),
    "right_margin": Mm(12.5),
    "top_margin": Mm(12.5),
    "bottom_margin": Mm(1.25),
    "image_width": Mm(30),
    "column_of_image_width": Mm(33),
    "image_height": Mm(30 / 16 * 9),
}

while True:
    imgs = glob(_IMGS)
    doc = Document()
    presentation = Presentation(_PPTX)

    style = doc.styles['Normal']
    font = style.font
    font.name = '游ゴシック Medium'

    slides = presentation.slides
    table = doc.add_table(rows=0, cols=2, style="Table Grid")

    if len(imgs) == len(slides):  # スライドの枚数と画像の枚数が一致するかを確認する
        for i in range(len(imgs)):
            slide = slides[i]
            img = imgs[i]
            cells = table.add_row().cells

            # 第一列にスライドのイメージを挿入する
            cells[0].BottomPadding = Mm(0)
            cells[0].LeftPadding = Mm(0)
            cells[0].RightPadding = Mm(0)
            cells[0].TopPadding = Mm(0)
            cells[0].paragraphs[0].add_run().add_picture(
                img,
                width=size_settings["image_width"],
                height=size_settings["image_height"]
            )

            # 第二列にスライドのコメント(ノート)を書き込む
            cells[1].text = slide.notes_slide.notes_text_frame.text
    else:
        print(len(imgs), len(slides), "Images does not exist as same count as slides.")
        raise

    # ページサイズ・余白・列幅を一括設定
    for section in doc.sections:
        section.page_height = size_settings["page_height"]
        section.page_width = size_settings["page_width"]
        section.left_margin = size_settings["left_margin"]
        section.right_margin = size_settings["right_margin"]
        section.top_margin = size_settings["top_margin"]
        section.bottom_margin = size_settings["bottom_margin"]

    table.width = doc.sections[0].page_width - (doc.sections[0].left_margin + doc.sections[0].right_margin)
    table.columns[0].width = size_settings["column_of_image_width"]
    table.columns[1].width = table.width - size_settings["column_of_image_width"]

    doc.save(_DOCX)
    _ = input("Do convert -> Press Enter.")
